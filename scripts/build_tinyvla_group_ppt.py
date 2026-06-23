from __future__ import annotations

from pathlib import Path
from zipfile import ZipFile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT_PATH = ROOT / "tinyVLA_group_meeting.pptx"
REF_PPT = ROOT / "0625-tinyVLA.pptx"
ASSET_DIR = ROOT / "tinyVLA_ppt_assets"
REF_ASSET_DIR = ASSET_DIR / "reference_style"

W, H = 13.33, 7.5
RED = RGBColor(198, 0, 56)
DARK_RED = RGBColor(156, 0, 42)
BLACK = RGBColor(20, 20, 20)
GRAY = RGBColor(86, 86, 86)
LIGHT_GRAY = RGBColor(246, 246, 246)
WHITE = RGBColor(255, 255, 255)


def extract_reference_assets() -> tuple[Path, Path]:
    REF_ASSET_DIR.mkdir(parents=True, exist_ok=True)
    campus = REF_ASSET_DIR / "campus.jpg"
    logo = REF_ASSET_DIR / "usstd_logo.png"
    if campus.exists() and logo.exists():
        return campus, logo

    with ZipFile(REF_PPT) as zf:
        campus.write_bytes(zf.read("ppt/media/image4.jpg"))
        logo.write_bytes(zf.read("ppt/media/image1.png"))
    return campus, logo


def add_picture_cover(slide, path: Path, x, y, w, h):
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def add_text(slide, text, x, y, w, h, size=24, color=BLACK, bold=False,
             font="Arial", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_rect(slide, x, y, w, h, fill=WHITE, line=None, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1.2)
    return shp


def add_line(slide, x1, y1, x2, y2, color=RED, width=2):
    c = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width)
    return c


def add_logo(slide, logo: Path):
    slide.shapes.add_picture(str(logo), Inches(10.2), Inches(0.28), width=Inches(2.55))


def add_header(slide, title: str, logo: Path):
    add_rect(slide, 0, 0.25, 6.15, 0.74, fill=RED, radius=True)
    add_text(slide, title, 0.42, 0.34, 5.2, 0.48, size=21, color=WHITE,
             bold=True, font="Microsoft YaHei", valign=MSO_ANCHOR.MIDDLE)
    add_logo(slide, logo)
    add_line(slide, 0.35, 1.2, 12.3, 1.2, color=RED, width=1.0)


def add_bullets(slide, bullets: list[str], x, y, w, h, size=18, line_spacing=1.18):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(9)
        p.line_spacing = line_spacing
        p.alignment = PP_ALIGN.LEFT
        r0 = p.add_run()
        r0.text = "➤  "
        r0.font.name = "Arial"
        r0.font.size = Pt(size)
        r0.font.bold = True
        r0.font.color.rgb = RED
        add_rich_runs(p, text, size)
    return box


def add_rich_runs(paragraph, text: str, size: int):
    parts = text.split("**")
    for idx, part in enumerate(parts):
        if not part:
            continue
        run = paragraph.add_run()
        run.text = part
        run.font.name = "Arial"
        run.font.size = Pt(size)
        run.font.bold = idx % 2 == 1
        run.font.color.rgb = RED if idx % 2 == 1 else BLACK


def add_content_frame(slide, x=0.35, y=1.35, w=12.3, h=5.85):
    add_rect(slide, x, y, w, h, fill=LIGHT_GRAY, line=RED, radius=True)


def add_caption(slide, text, x, y, w):
    add_text(slide, text, x, y, w, 0.25, size=10.5, color=GRAY, font="Microsoft YaHei",
             align=PP_ALIGN.CENTER)


def add_image_fit(slide, path: Path, x, y, w, h, caption=""):
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    if caption:
        add_caption(slide, caption, x, y + h + 0.08, w)


def cover(prs, campus: Path, logo: Path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_picture_cover(slide, campus, 0, 0, W, 3.8)
    oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.05), Inches(2.35), Inches(15.4), Inches(3.0))
    oval.fill.solid()
    oval.fill.fore_color.rgb = WHITE
    oval.line.fill.background()
    slide.shapes.add_picture(str(logo), Inches(9.15), Inches(0.12), width=Inches(3.75))
    title = "TinyVLA：快速、数据高效的 VLA 模型"
    add_text(slide, title, 1.35, 3.95, 10.65, 0.75, size=31, color=RED, bold=True,
             font="Microsoft YaHei", align=PP_ALIGN.CENTER)
    add_line(slide, 1.35, 5.35, 11.95, 5.35, color=RED, width=0.9)
    add_text(slide, "IEEE ROBOTICS AND AUTOMATION LETTERS", 2.6, 5.65, 8.2, 0.45,
             size=19, color=BLACK, font="Arial", align=PP_ALIGN.CENTER)
    add_rect(slide, 4.1, 6.23, 2.25, 0.52, fill=RED, radius=True)
    add_text(slide, "汇报人：组会汇报", 4.22, 6.29, 2.0, 0.35, size=15, color=WHITE,
             bold=True, font="Microsoft YaHei", align=PP_ALIGN.CENTER)
    add_rect(slide, 7.0, 6.23, 1.95, 0.52, fill=RED, radius=True)
    add_text(slide, "2026-06-23", 7.13, 6.29, 1.7, 0.35, size=15, color=WHITE,
             bold=True, font="Arial", align=PP_ALIGN.CENTER)


def toc(prs, campus: Path, logo: Path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_picture_cover(slide, campus, 0, 0, 5.55, H)
    add_rect(slide, 5.55, 0, 0.12, H, fill=RED)
    add_rect(slide, 5.05, 2.35, 1.08, 2.85, fill=RED, radius=True)
    add_text(slide, "目录", 5.32, 3.2, 0.5, 1.0, size=30, color=WHITE, bold=True,
             font="Microsoft YaHei", align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_logo(slide, logo)
    items = [("01", "摘要", "Abstract"), ("02", "方法", "Method"),
             ("03", "实验", "Experiment"), ("04", "结论", "Conclusion")]
    y = 1.55
    for num, cn, en in items:
        add_text(slide, num, 6.9, y, 1.0, 0.55, size=40, color=RED, bold=True, font="Arial")
        add_text(slide, cn, 8.05, y + 0.02, 2.2, 0.35, size=25, color=BLACK,
                 font="Microsoft YaHei")
        add_text(slide, en, 8.1, y + 0.45, 2.2, 0.35, size=18, color=BLACK, font="Arial")
        y += 1.55


def section(prs, campus: Path, idx: int, cn: str, en: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_picture_cover(slide, campus, 0, 0, W, 4.35)
    add_rect(slide, 0, 4.38, W, 0.12, fill=RED)
    add_rect(slide, 4.65, 3.75, 4.0, 1.05, fill=RED, radius=True)
    add_text(slide, f"PART {idx:02d}", 5.25, 4.0, 2.8, 0.55, size=38, color=WHITE,
             bold=True, font="Arial", align=PP_ALIGN.CENTER)
    add_text(slide, cn, 0.0, 5.35, W, 0.55, size=38, color=GRAY, bold=True,
             font="Microsoft YaHei", align=PP_ALIGN.CENTER)
    add_text(slide, en, 0.0, 6.15, W, 0.45, size=22, color=BLACK,
             font="Arial", align=PP_ALIGN.CENTER)
    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(6.45), Inches(7.18), Inches(0.42), Inches(0.34))
    tri.fill.solid()
    tri.fill.fore_color.rgb = RED
    tri.line.fill.background()


def content_slide(prs, logo: Path, title: str, bullets: list[str], image: Path | None = None,
                  caption: str = "", image_w=4.9, image_h=3.15):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, title, logo)
    add_content_frame(slide)
    if image is None:
        add_bullets(slide, bullets, 0.85, 1.72, 11.35, 4.9, size=21, line_spacing=1.22)
    else:
        add_bullets(slide, bullets, 0.78, 1.68, 6.2, 5.0, size=17, line_spacing=1.12)
        add_image_fit(slide, image, 7.2, 2.05, image_w, image_h, caption=caption)
    return slide


def two_image_slide(prs, logo: Path, title: str, bullets: list[str],
                    left_img: Path, left_caption: str, right_img: Path, right_caption: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, title, logo)
    add_content_frame(slide)
    add_bullets(slide, bullets, 0.78, 1.65, 11.65, 1.75, size=15.5, line_spacing=1.05)
    add_line(slide, 0.7, 3.62, 12.25, 3.62, color=RGBColor(210, 210, 210), width=0.8)
    add_image_fit(slide, left_img, 0.95, 3.9, 5.65, 2.18, caption=left_caption)
    add_image_fit(slide, right_img, 7.05, 3.9, 5.1, 2.18, caption=right_caption)


def build():
    campus, logo = extract_reference_assets()
    f = ASSET_DIR
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)

    cover(prs, campus, logo)
    toc(prs, campus, logo)

    section(prs, campus, 1, "摘要", "Abstract")
    content_slide(
        prs, logo, "Abstract",
        [
            "TinyVLA targets a practical VLA problem: current models are **slow at inference** and often require large-scale robot pre-training.",
            "The proposed family of compact VLA models improves **inference speed** and **data efficiency**, avoiding an OpenX-style robot pre-training stage.",
            "The key idea is to initialize the policy with a compact pre-trained VLM and attach a **diffusion policy decoder** for continuous robot actions.",
        ],
        f / "fig1_latency_success.png", "Inference latency vs. average success rate", 4.65, 3.35
    )
    content_slide(
        prs, logo, "Motivation",
        [
            "Large VLA models such as OpenVLA inherit strong vision-language priors, but their 7B-scale VLM backbone and autoregressive action tokens make real-time control expensive.",
            "TinyVLA asks whether a smaller VLM can preserve semantic generalization while a non-autoregressive policy head handles precise action generation.",
            "**Takeaway:** robot policy design should optimize both the multimodal backbone and the action decoder.",
        ]
    )

    section(prs, campus, 2, "方法", "Method")
    content_slide(
        prs, logo, "TinyVLA Architecture",
        [
            "VLM backbone: TinyVLA trains compact VLMs with **70M-1.4B parameters**, using **Pythia** as the language-model backend.",
            "Training recipe: the VLM family follows the **LLaVA** training pipeline and keeps the visual backbone plus vision-language alignment module during robot fine-tuning.",
            "Fine-tuning: LoRA is inserted into Transformer attention Q/K/V while most pre-trained weights are frozen, preserving language and visual priors.",
            "**Core innovation:** replace autoregressive action-token prediction with a diffusion policy decoder.",
        ],
        f / "fig2_architecture.png", "TinyVLA model architecture", 4.85, 3.05
    )
    content_slide(
        prs, logo, "Action Model: Diffusion Policy Decoder",
        [
            "Instead of discretizing continuous actions into tokens, TinyVLA uses **Diffusion Policy (DP)** as the policy head.",
            "DP formulates action generation with **Denoising Diffusion Probabilistic Models (DDPMs)**: the network predicts denoising noise and iteratively recovers an action trajectory.",
            "Pipeline: VLM encodes observations and language; adaptive pooling + LayerNorm produces compact features; proprioception is concatenated and passed through a 3-layer MLP as DP conditioning.",
            "This decoder directly outputs continuous high-dimensional robot actions and avoids repeated next-token inference.",
        ]
    )

    section(prs, campus, 3, "实验", "Experiment")
    content_slide(
        prs, logo, "Experimental Setup",
        [
            "Simulation: MetaWorld 50 tasks, grouped by difficulty; each task uses 50 demonstrations and is evaluated over three seeds.",
            "Real robots: Franka single-arm tasks and bimanual UR5 cooperative tasks test both manipulation diversity and embodiment transfer.",
            "Generalization: instruction, view, background, distractor, illumination, appearance and spatial shifts.",
        ],
        f / "fig3_robot_setup.png", "Real-robot setups", 5.15, 2.65
    )
    two_image_slide(
        prs, logo, "Multi-task and Real-world Results",
        [
            "In simulation, TinyVLA-H reaches an average success rate of 31.6%, compared with 10.5% for Diffusion Policy.",
            "On real single-arm tasks, TinyVLA-H achieves **94.0%** average success, compared with **68.3%** for OpenVLA.",
            "On bimanual UR5 tasks, OpenVLA nearly fails, while TinyVLA-H still reaches 44.5% average success.",
        ],
        f / "table2_real_world.png", "Real-world single-arm results",
        f / "table4_latency.png", "Inference latency on A6000",
    )
    two_image_slide(
        prs, logo, "Generalization Evaluation",
        [
            "TinyVLA handles unseen colors, objects and object-function compositions, indicating that VLM semantic priors transfer into policy learning.",
            "Under camera-view changes, TinyVLA remains more robust than Diffusion Policy and often approaches or exceeds OpenVLA.",
            "Spatial tests move targets outside the training area; TinyVLA still solves part of the position-sensitive tasks.",
        ],
        f / "fig5_view_generalization.png", "View generalization",
        f / "fig9_spatial_generalization.png", "Spatial generalization",
    )
    content_slide(
        prs, logo, "Ablation Study",
        [
            "Model scale matters: TinyVLA-0.4B fails more often due to misinterpreted instructions and localization errors; 1.3B substantially reduces these failures.",
            "TinyVLA-3B uses the pre-trained **PaliGemma** model in the ablation, suggesting stronger localization-oriented VLMs can further help manipulation.",
            "Policy-head choice matters: the diffusion policy head outperforms MLP and ACT alternatives across the five real-robot tasks.",
            "Speed comes from both sides: a compact VLM lowers encoding cost, while DP avoids autoregressive action-token generation.",
        ],
        f / "fig10_failure_types.png", "Failure types across VLM sizes", 4.4, 3.15
    )

    section(prs, campus, 4, "结论", "Conclusion")
    content_slide(
        prs, logo, "Conclusion",
        [
            "TinyVLA demonstrates that strong VLA policies do not necessarily require large-scale robot pre-training.",
            "The most important design is the division of labor: a compact VLM provides semantic understanding, while a diffusion policy decoder generates continuous actions.",
            "Compared with OpenVLA, TinyVLA-H is more data-efficient, uses far fewer parameters and achieves much lower latency.",
            "For bimanual service tasks, the paper suggests a useful direction: combine compact VLA backbones with task-stage memory and arm-role-aware action modeling.",
        ]
    )
    content_slide(
        prs, logo, "Discussion",
        [
            "What exactly transfers from the VLM: object semantics, spatial grounding, instruction following, or all of them?",
            "How far can diffusion policy scale when the action horizon grows or the task requires long-horizon bimanual coordination?",
            "For beverage-service robots, TinyVLA motivates a compact VLA policy, but task phases and left/right arm roles may need explicit structure.",
        ]
    )

    prs.save(OUTPUT_PATH)
    print(OUTPUT_PATH)


if __name__ == "__main__":
    build()
